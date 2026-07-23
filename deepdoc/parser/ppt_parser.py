#
#  Copyright 2025 The InfiniFlow Authors. All Rights Reserved.
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#

"""
PowerPoint 文件解析器

基于 python-pptx 库解析 .pptx 文件，提取每页幻灯片中的文本内容。
支持处理多种形状类型：文本框、表格、组合形状，并按位置排序保证阅读顺序。
"""

import logging
from io import BytesIO
from pptx import Presentation


class RAGFlowPptParser:
    """PowerPoint 演示文稿解析器

    解析 PPTX 文件的每一页幻灯片，提取其中的文本内容。
    支持带格式文本（项目符号按层级缩进）、表格、组合形状的处理。
    形状按（top位置, left位置）排序以保证近似阅读顺序。
    """

    def __init__(self):
        super().__init__()
        # 形状排序缓存：避免对同一组形状重复排序
        self._shape_cache = {}

    def __sort_shapes(self, shapes):
        """按位置对形状排序

        排序规则：先按 top（纵坐标）每 10 像素为一组，同组内按 left（横坐标）排序。
        这样保证了从上到下、从左到右的近似阅读顺序。
        使用缓存避免对同一组形状重复排序。

        Args:
            shapes: PPT 形状列表

        Returns:
            排序后的形状列表
        """
        cache_key = id(shapes)
        if cache_key not in self._shape_cache:
            self._shape_cache[cache_key] = sorted(
                shapes,
                key=lambda x: ((x.top if x.top is not None else 0) // 10, x.left if x.left is not None else 0)
            )
        return self._shape_cache[cache_key]

    def __get_bulleted_text(self, paragraph):
        """获取带项目符号格式的段落文本

        检测段落是否为项目符号（bullet）类型，如果是则按层级缩进添加前缀。

        Args:
            paragraph: python-pptx 的 Paragraph 对象

        Returns:
            带格式的文本字符串，项目符号段落前缀类似 "  *"、"    *" 等
        """
        # 检测三种项目符号类型：字符符号(buChar)、自动编号(buAutoNum)、图片符号(buBlip)
        is_bulleted = bool(paragraph._p.xpath("./a:pPr/a:buChar")) or bool(paragraph._p.xpath("./a:pPr/a:buAutoNum")) or bool(paragraph._p.xpath("./a:pPr/a:buBlip"))
        if is_bulleted:
            return f"{'  '* paragraph.level}.{paragraph.text}"
        else:
            return paragraph.text

    def __extract(self, shape):
        """从单个形状中提取文本内容

        根据形状类型分别处理：
        - 文本框：提取所有段落文本
        - 表格（shape_type == 19）：以 "表头: 值" 格式输出
        - 组合形状（shape_type == 6）：递归提取子形状
        - 其他类型：尝试获取 text 属性

        Args:
            shape: python-pptx 的 Shape 对象

        Returns:
            提取的文本字符串，失败返回空字符串
        """
        try:
            # 优先尝试提取文本框内容
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                text_frame = shape.text_frame
                texts = []
                for paragraph in text_frame.paragraphs:
                    if paragraph.text.strip():
                        texts.append(self.__get_bulleted_text(paragraph))
                return "\n".join(texts)

            # 安全获取形状类型
            try:
                shape_type = shape.shape_type
            except NotImplementedError:
                if hasattr(shape, 'text'):
                    return shape.text.strip()
                return ""

            # 处理表格（MSO_SHAPE_TYPE.TABLE = 19）
            if shape_type == 19:
                tb = shape.table
                rows = []
                for i in range(1, len(tb.rows)):
                    rows.append("; ".join([tb.cell(
                        0, j).text + ": " + tb.cell(i, j).text for j in range(len(tb.columns)) if tb.cell(i, j)]))
                return "\n".join(rows)

            # 处理组合形状（MSO_SHAPE_TYPE.GROUP = 6）
            if shape_type == 6:
                texts = []
                for p in self.__sort_shapes(shape.shapes):
                    t = self.__extract(p)
                    if t:
                        texts.append(t)
                return "\n".join(texts)

            return ""

        except Exception as e:
            logging.error(f"Error processing shape: {str(e)}")
            return ""

    def __call__(self, fnm, from_page, to_page, callback=None):
        """解析 PPT 文件

        Args:
            fnm: PPTX 文件路径或二进制内容
            from_page: 起始页码（从 0 开始）
            to_page: 结束页码（不包含）
            callback: 进度回调函数

        Returns:
            字符串列表，每项为一页幻灯片的文本内容
        """
        ppt = Presentation(fnm) if isinstance(
            fnm, str) else Presentation(
            BytesIO(fnm))
        txts = []
        self.total_page = len(ppt.slides)
        for i, slide in enumerate(ppt.slides):
            if i < from_page:
                continue
            if i >= to_page:
                break
            texts = []
            # 按位置排序后逐个提取形状文本
            for shape in self.__sort_shapes(slide.shapes):
                txt = self.__extract(shape)
                if txt:
                    texts.append(txt)
            txts.append("\n".join(texts))

        return txts
